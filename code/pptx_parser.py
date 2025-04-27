# -*- coding: utf-8 -*-
import json
import pandas as pd
from pptx import Presentation
import os
from typing import Dict, Any, List
from datetime import datetime

class PPTParser:
    """用于解析 PowerPoint 文件并将其转换为 JSON 格式的类，每页视为一个 sheet。"""

    def __init__(self, file_path: str, doc_type: str = "ppt"):
        """
        初始化 PPTParser 类。
        :param file_path: PPT 文件路径
        :param doc_type: 文档类型
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPT file not found: {file_path}")

        self.file_path = file_path
        self.file_name = os.path.basename(file_path)
        self.doc_type = doc_type
        self.result = {
            "doc_type": self.doc_type,
            "file_name": self.file_name,
            "tables": []
        }

    def get_slide_title(self, slide: 'pptx.slide.Slide') -> str:
        """
        获取幻灯片的大标题，若无标题则返回默认名称。
        :param slide: PPT 幻灯片对象
        :return: 标题字符串
        """
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                text = shape.text.strip()
                if text:  # 优先返回第一个非空 placeholder 文本作为标题
                    return text
        return "Untitled"  # 默认标题

    def extract_table_data(self, table: 'pptx.table.Table') -> List[List[str]]:
        """
        提取表格数据，处理合并单元格。
        :param table: PPT 表格对象
        :return: 二维数据列表
        """
        rows = len(table.rows)
        cols = len(table.columns)
        data = [["" for _ in range(cols)] for _ in range(rows)]

        # 遍历所有单元格
        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                # 如果单元格未被填充（非合并区域）
                if not data[row_idx][col_idx]:
                    value = cell.text.strip()
                    if isinstance(value, datetime):
                        value = value.isoformat()
                    elif value is not None:
                        value = str(value)

                    # 处理合并单元格
                    if cell.is_merge_origin:
                        # 填充合并区域
                        try:
                            # 推断合并范围（基于 is_spanned）
                            merge_rows = 1
                            merge_cols = 1
                            for r in range(row_idx + 1, rows):
                                next_cell = table.cell(r, col_idx)
                                if next_cell.is_spanned:
                                    merge_rows += 1
                                else:
                                    break
                            for c in range(col_idx + 1, cols):
                                next_cell = table.cell(row_idx, c)
                                if next_cell.is_spanned:
                                    merge_cols += 1
                                else:
                                    break
                            # 填充合并区域
                            for r in range(row_idx, row_idx + merge_rows):
                                for c in range(col_idx, col_idx + merge_cols):
                                    if r < rows and c < cols:
                                        data[r][c] = value
                        except AttributeError:
                            # 如果 is_spanned 不可用，直接填充当前单元格
                            data[row_idx][col_idx] = value
                    else:
                        # 非合并单元格，直接填充
                        data[row_idx][col_idx] = value

        return data

    def parse(self) -> Dict[str, Any]:
        """
        解析 PPT 文件，提取每页的表格和文本内容，每页视为一个 sheet。
        :return: 解析后的 JSON 数据，包含 doc_type、file_name 和 tables
        """
        try:
            prs = Presentation(self.file_path)
            self.result["tables"] = []

            for slide_idx, slide in enumerate(prs.slides):
                # 获取 sheet 名称（大标题）
                sheet_name = self.get_slide_title(slide)

                # 提取文本内容（非表格部分）
                text_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame and not shape.has_table:
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)

                # 提取表格内容
                table_data = None
                for shape in slide.shapes:
                    if shape.has_table:
                        table_data = self.extract_table_data(shape.table)
                        break  # 假设每页最多一个表格

                # 构造表格数据（如果存在）
                if table_data:
                    # 过滤空白行
                    table_data = [row for row in table_data if any(cell != "" for cell in row)]
                    if table_data:
                        # 假设第一行为表头
                        headers = table_data[0]
                        header_counts = {}
                        new_headers = []
                        for i, header in enumerate(headers):
                            header = str(header).strip() if header else f"Column_{i+1}"
                            if header in header_counts:
                                header_counts[header] += 1
                                new_headers.append(f"{header}_{header_counts[header]}")
                            else:
                                header_counts[header] = 0
                                new_headers.append(header)

                        df_data = table_data[1:] if len(table_data) > 1 else []
                        df = pd.DataFrame(df_data, columns=new_headers)

                        csv_content = df.to_csv(index=False, encoding='utf-8')
                        rows = df.to_dict(orient="records")
                    else:
                        csv_content = ""
                        rows = []
                else:
                    csv_content = ""
                    rows = []

                # 添加到结果
                self.result["tables"].append({
                    "sheet": sheet_name,
                    "data": csv_content,
                    "rows": rows,
                    "text": "\n".join(text_parts)  # 非表格文本内容
                })

        except Exception as e:
            raise Exception(f"Error processing PPT file: {str(e)}")

        return self.result

    def save_sheets_to_files(self, output_dir: str = "output"):
        """
        将每页的数据保存为单独的 JSON 文件，文件名为“ppt名称_sheet名称_2.json”。
        为重复的 sheet 名称添加数字后缀（项目概况, 项目概况1, 项目概况2 等）。
        :param output_dir: 输出目录
        """
        # 确保输出目录存在
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # 跟踪 sheet 名称使用次数
        sheet_name_counts = {}
        for table in self.result["tables"]:
            sheet_name = table["sheet"]
            # 替换非法文件名字符
            sheet_name = sheet_name.replace("/", "_").replace("\\", "_").replace(":", "_")

            # 处理重复 sheet 名称
            if sheet_name in sheet_name_counts:
                sheet_name_counts[sheet_name] += 1
                unique_sheet_name = f"{sheet_name}{sheet_name_counts[sheet_name]}"
            else:
                sheet_name_counts[sheet_name] = 0
                unique_sheet_name = sheet_name

            # 构造输出文件名：ppt名称_sheet名称_0.json
            base_name = os.path.splitext(self.file_name)[0]  # 去掉扩展名
            output_file_name = f"{base_name}_{unique_sheet_name}_2.json"
            output_path = os.path.join(output_dir, output_file_name)

            # 构造单个 sheet 的 JSON 数据
            sheet_json = {
                "doc_type": self.doc_type,
                "file_name": self.file_name,
                "tables": [table]
            }

            # 保存到文件
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(sheet_json, f, ensure_ascii=False, indent=2)
            print(f"Saved sheet '{unique_sheet_name}' to: {output_path}")

def main():
    # PPT 文件路径
    file_path = r"C:\Users\dreame\Desktop\电子元件RAG\ppt\ERP优化总结方案.pptx"

    # 验证文件存在性
    print(f"Checking file: {file_path}")
    if not os.path.exists(file_path):
        print(f"Error: File does not exist at {file_path}")
        return

    # 使用 PPTParser 解析 PPT 文件
    try:
        parser = PPTParser(file_path)
        result = parser.parse()

        # 打印解析结果
        print("\n=== Parsed PPT ===")
        print(f"Doc Type: {result['doc_type']}")
        print(f"File Name: {result['file_name']}")
        print(f"Total Slides: {len(result['tables'])}")
        for i, table in enumerate(result['tables']):
            print(f"\nSlide {i + 1} - Sheet Name: {table['sheet']}")
            print(f"Number of Table Rows: {len(table['rows'])}")
            print(f"Table Data Length: {len(table['data'])} characters")
            print(f"Text Content:\n{table['text']}")

        # 保存每页为单独的 JSON 文件
        parser.save_sheets_to_files(output_dir="llm_output_test")

    except Exception as e:
        print(f"Error: {str(e)}")

if __name__ == "__main__":
    main()